from typing import Dict, List, Any
import os
from pathlib import Path
from pptx import Presentation
from docx import Document
from .document_processor import DocumentProcessor

class OfficeProcessor(DocumentProcessor):
    """Processor for Microsoft Office documents (PPTX, DOCX)"""
    
    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process Office document and return extracted text with metadata"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext in ['.pptx', '.ppt']:
                return self._process_presentation(file_path)
            elif ext in ['.docx', '.doc']:
                return self._process_document(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
                
        except Exception as e:
            error_msg = str(e)
            if "file format" in error_msg.lower() or "not a word file" in error_msg.lower():
                raise ValueError("This appears to be an unsupported Word format. Please use .docx format.")
            elif "not a powerpoint" in error_msg.lower():
                raise ValueError("This appears to be an unsupported PowerPoint format. Please use .pptx format.")
            
            print(f"Error processing office document: {error_msg}")
            return {
                'texts': ['Error processing office document'],
                'metadata': [{
                    'source': str(file_path),
                    'type': 'error',
                    'error': error_msg
                }]
            }

    def _process_presentation(self, file_path: str) -> Dict[str, List[Any]]:
        """Process PPTX files"""
        try:
            prs = Presentation(file_path)
            all_content = []
            all_metadata = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"[Slide {slide_num}]")
                
                if slide.shapes.title:
                    slide_content.append(f"Title: {slide.shapes.title.text}")
                
                # Process each shape in the slide
                for shape in slide.shapes:
                    # Extract text content
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # Process tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        table_content = []
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                table_content.append(row_text)
                        if table_content:
                            slide_content.append(
                                "Table content:\n" + "\n".join(table_content)
                            )
                    
                    # Mark presence of images
                    if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture
                        slide_content.append("[Image in slide]")
                
                content = "\n".join(slide_content)
                if content.strip():
                    all_content.append(content)
                    all_metadata.append({
                        'source': str(file_path),
                        'type': 'presentation',
                        'slide': slide_num,
                        'total_slides': len(prs.slides)
                    })
            
            # Handle empty presentation because people are stoopid
            if not all_content:
                return {
                    'texts': ['Empty presentation'],
                    'metadata': [{
                        'source': str(file_path),
                        'type': 'presentation',
                        'total_slides': len(prs.slides)
                    }]
                }
            
            return {
                'texts': all_content,
                'metadata': all_metadata
            }
        except Exception as e:
            if 'not a PowerPoint file' in str(e):
                raise ValueError("This appears to be an unsupported PowerPoint format. Please convert to .pptx format.")
            raise

    def _process_document(self, file_path: str) -> Dict[str, List[Any]]:
        """Process DOCX files"""
        try:
            doc = Document(file_path)
            all_content = []
            all_metadata = []
            section_content = []
            section_num = 1
            total_sections = max(1, len(doc.paragraphs) // 10 + 1)  # Estimate sections
            
            # Process paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    section_content.append(para.text.strip())
                
                # Split into chunks if section is getting too large
                if len("\n".join(section_content)) > 1000:
                    all_content.append("\n".join(section_content))
                    all_metadata.append({
                        'source': str(file_path),
                        'type': 'document',
                        'section': section_num,
                        'total_sections': total_sections,
                        'has_tables': len(doc.tables) > 0
                    })
                    section_content = []
                    section_num += 1
            
            # Process tables, bhai be matlab ki mehnat
            for table_num, table in enumerate(doc.tables, 1):
                table_content = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_content.append(row_text)
                if table_content:
                    section_content.append(
                        f"Table {table_num} content:\n" + "\n".join(table_content)
                    )
                
                if len("\n".join(section_content)) > 1000:
                    all_content.append("\n".join(section_content))
                    all_metadata.append({
                        'source': str(file_path),
                        'type': 'document',
                        'section': section_num,
                        'total_sections': total_sections,
                        'has_tables': len(doc.tables) > 0
                    })
                    section_content = []
                    section_num += 1
            
            # Add any remaining content
            if section_content:
                all_content.append("\n".join(section_content))
                all_metadata.append({
                    'source': str(file_path),
                    'type': 'document',
                    'section': section_num,
                    'total_sections': total_sections,
                    'has_tables': len(doc.tables) > 0
                })
            
            # Handle empty document
            if not all_content:
                return {
                    'texts': ['Empty document'],
                    'metadata': [{
                        'source': str(file_path),
                        'type': 'document',
                        'total_sections': 1,
                        'has_tables': len(doc.tables) > 0
                    }]
                }
            
            return {
                'texts': all_content,
                'metadata': all_metadata
            }
        except Exception as e:
            if 'not a Word file' in str(e):
                raise ValueError("This appears to be an unsupported Word format. Please convert to .docx format.")
            raise