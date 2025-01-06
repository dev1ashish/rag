import pandas as pd
from PyPDF2 import PdfReader
import openai
from typing import Dict, List, Any, Optional
import os
from PIL import Image
import base64
import requests
from bs4 import BeautifulSoup
import io
from pathlib import Path
from pptx import Presentation
from docx import Document
from utils.data_visualization import DataVisualizationManager

class DocumentProcessor:
    """Base class for document processors"""
    def process(self, file_path: str) -> Dict[str, List[Any]]:
        raise NotImplementedError

class ExcelProcessor(DocumentProcessor):
    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process Excel file and return extracted text with metadata."""
        try:
            # Read the Excel file
            df = pd.read_excel(file_path)
            
            # Store DataFrame for visualization
            DataVisualizationManager.store_dataframe(df, os.path.basename(file_path))
            
            # Create enhanced metadata
            structured_metadata = {
                'source': str(file_path),
                'type': 'excel_structured',
                'total_rows': len(df),
                'total_columns': len(df.columns),
                'columns': ', '.join(df.columns.tolist()),
                'column_types': {str(col): str(dtype) for col, dtype in df.dtypes.items()},
                'has_numerical_data': any(df[col].dtype.kind in 'biufc' for col in df.columns),
                'numerical_columns': [col for col in df.columns if df[col].dtype.kind in 'biufc'],
                'categorical_columns': [col for col in df.columns if df[col].dtype.kind not in 'biufc']
            }

            # Add basic statistics for numerical columns
            if structured_metadata['has_numerical_data']:
                stats = {}
                for col in structured_metadata['numerical_columns']:
                    try:
                        stats[col] = {
                            'min': float(df[col].min()),
                            'max': float(df[col].max()),
                            'mean': float(df[col].mean()),
                            'median': float(df[col].median())
                        }
                    except Exception:
                        continue
                structured_metadata['numerical_stats'] = stats

            # Convert to string representation for text embedding
            text_content = [df.to_string(index=False)]
            
            return {
                'texts': text_content,
                'metadata': [structured_metadata],
                'dataframe': df  # Include DataFrame in return
            }
        except Exception as e:
            print(f"Error processing Excel file: {str(e)}")
            return {
                'texts': ['Error processing Excel file'],
                'metadata': [{
                    'source': str(file_path), 
                    'type': 'error',
                    'error': str(e)
                }]
            }

class PDFProcessor(DocumentProcessor):
    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process PDF file and return extracted text with metadata."""
        try:
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                texts = []
                metadata = []
                
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        texts.append(text)
                        metadata.append({
                            'source': os.path.basename(file_path),
                            'type': 'pdf',
                            'page': page_num + 1
                        })
                
                return {
                    'texts': texts,
                    'metadata': metadata
                }
        except Exception as e:
            print(f"Error processing PDF file: {str(e)}")
            return {
                'texts': ['Error processing PDF file'],
                'metadata': [{'source': str(file_path), 'type': 'error'}]
            }

class AudioProcessor(DocumentProcessor):
    def __init__(self, api_key: str):
        """Initialize audio processor with OpenAI API key."""
        self.api_key = api_key
        openai.api_key = api_key

    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process audio file and return transcription with metadata."""
        try:
            with open(file_path, "rb") as audio_file:
                transcript = openai.Audio.transcribe(
                    model="whisper-1",
                    file=audio_file
                )

            text = transcript.get("text") if isinstance(transcript, dict) else str(transcript)

            return {
                'texts': [text],
                'metadata': [{
                    'source': str(file_path),
                    'type': 'audio',
                    'model': 'whisper-1'
                }]
            }
        except Exception as e:
            print(f"Error processing audio file: {str(e)}")
            return {
                'texts': ['Error processing audio file'],
                'metadata': [{'source': str(file_path), 'type': 'error'}]
            }

class ImageProcessor(DocumentProcessor):
    def __init__(self, api_key: str):
        self.api_key = api_key
        openai.api_key = api_key

    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process image file and return description with metadata."""
        try:
            image = Image.open(file_path)
            
            if image.mode in ('RGBA', 'P'):
                image = image.convert('RGB')
            
            max_size = (1024, 1024)
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
            
            buffered = io.BytesIO()
            image.save(buffered, format=image.format if image.format else "JPEG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            
            print(f"Attempting to process image: {os.path.basename(file_path)}")
            
            response = openai.ChatCompletion.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this image in detail, including any text, objects, and important information visible."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{img_str}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=300
            )
            
            description = response.choices[0].message['content']
            
            return {
                'texts': [description],
                'metadata': [{
                    'source': str(file_path),
                    'type': 'image',
                    'format': image.format,
                    'size': image.size
                }]
            }
        except Exception as e:
            print(f"Error processing image file: {str(e)}")
            return {
                'texts': ['Error processing image file'],
                'metadata': [{'source': str(file_path), 'type': 'error', 'error': str(e)}]
            }

class WebProcessor(DocumentProcessor):
    def process(self, url: str) -> Dict[str, List[Any]]:
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            
            try:
                soup = BeautifulSoup(response.text, 'html.parser')
            except Exception as e:
                return {
                    'texts': ['Error parsing webpage'],
                    'metadata': [{'source': str(url), 'type': 'error'}]
                }
            
            for element in soup(['script', 'style', 'header', 'footer', 'nav']):
                element.decompose()
            
            text = soup.get_text(separator='\n')
            
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            
            return {
                'texts': [text],
                'metadata': [{
                    'source': url,
                    'type': 'web',
                    'title': soup.title.string if soup.title else url
                }]
            }
        except Exception as e:
            print(f"Error processing webpage: {str(e)}")
            return {
                'texts': ['Error processing webpage'],
                'metadata': [{'source': str(url), 'type': 'error'}]
            }

class OfficeProcessor(DocumentProcessor):
    def process(self, file_path: str) -> Dict[str, List[Any]]:
        """Process Office document and return extracted text with metadata"""
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext in ['.pptx', '.ppt']:
                return self._process_presentation(file_path)
            elif ext in ['.docx', '.doc']:
                return self._process_document(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
                
        except Exception as e:
            print(f"Error processing office document: {str(e)}")
            return {
                'texts': ['Error processing office document'],
                'metadata': [{
                    'source': str(file_path),
                    'type': 'error',
                    'error': str(e)
                }]
            }

    def _process_presentation(self, file_path: str) -> Dict[str, List[Any]]:
        """Process PPTX files"""
        try:
            prs = Presentation(file_path)
            all_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"[Slide {slide_num}]")
                
                if slide.shapes.title:
                    slide_content.append(f"Title: {slide.shapes.title.text}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    if hasattr(shape, "has_table") and shape.has_table:
                        table_content = []
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                table_content.append(row_text)
                        if table_content:
                            slide_content.append(
                                "Table content:\n" + "\n".join(table_content)
                            )
                
                content = "\n".join(slide_content)
                if content.strip():
                    all_content.append(content)
            
            return {
                'texts': all_content,
                'metadata': [{
                    'source': str(file_path),
                    'type': 'presentation',
                    'total_slides': len(prs.slides)
                }]
            }
        except Exception as e:
            if 'not a PowerPoint file' in str(e):
                raise ValueError("This appears to be an unsupported PowerPoint format. Please convert to .pptx format.")
            raise

    def _process_document(self, file_path: str) -> Dict[str, List[Any]]:
        """Process DOCX files"""
        try:
            doc = Document(file_path)
            all_content = []
            section_content = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    section_content.append(para.text.strip())
                
                if len("\n".join(section_content)) > 1000:
                    all_content.append("\n".join(section_content))
                    section_content = []
            
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_content.append(row_text)
                if table_content:
                    section_content.append(
                        "Table content:\n" + "\n".join(table_content)
                    )
                    
                if len("\n".join(section_content)) > 1000:
                    all_content.append("\n".join(section_content))
                    section_content = []
            
            if section_content:
                all_content.append("\n".join(section_content))
            
            return {
                'texts': all_content,
                'metadata': [{
                    'source': str(file_path),
                    'type': 'document',
                    'total_paragraphs': len(doc.paragraphs),
                    'total_tables': len(doc.tables)
                }]
            }
        except Exception as e:
            if 'not a Word file' in str(e):
                raise ValueError("This appears to be an unsupported Word format. Please convert to .docx format.")
            raise

def get_processor(file_type: str, api_key: Optional[str] = None) -> DocumentProcessor:
    """Factory function to get appropriate processor"""
    processor_mapping = {
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pdf': 'pdf',
        '.mp3': 'audio',
        '.wav': 'audio',
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.pptx': 'office',
        '.ppt': 'office',
        '.docx': 'office',
        '.doc': 'office'
    }
    
    if file_type.startswith('.'):
        processor_type = processor_mapping.get(file_type.lower())
    else:
        processor_type = file_type if file_type in {'excel', 'pdf', 'audio', 'image', 'web', 'office'} else None
    
    if processor_type is None:
        raise ValueError(f"Unsupported file type: {file_type}")
        
    if processor_type in ['audio', 'image'] and api_key is None:
        raise ValueError(f"API key required for {processor_type} processing")
    
    processors = {
        'excel': lambda: ExcelProcessor(),
        'pdf': lambda: PDFProcessor(),
        'audio': lambda: AudioProcessor(api_key),
        'image': lambda: ImageProcessor(api_key),
        'web': lambda: WebProcessor(),
        'office': lambda: OfficeProcessor()
    }
    
    return processors[processor_type]()